# -*- coding: utf-8 -*-
"""
Simple_IMChat 软件体系结构课程设计 PPT 自动生成脚本 (高级美化与内容扩充版 v4)
作者: toh / Antigravity AI
说明:
    1. 采用 python-pptx 物理构建幻灯片，页面比例为 16:9。
    2. 主题选用高级极夜蓝 (RGB: 22, 26, 30) 作为背景，暗灰色 (RGB: 38, 42, 46) 为卡片色。
    3. 亮天蓝 (RGB: 0, 210, 255) 作为主强调色，高亮青色 (RGB: 0, 255, 180) 作为标志与饰条色。
    4. 扩充至 13 页：将 4+1 视图完全拆解为独立页面（场景、逻辑、过程-确认、过程-路由、开发、物理部署），增加各章细节参数与架构决策描述。
    5. 每页 Slide 均内嵌现场汇报演讲备注稿。
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    # 初始化幻灯片
    prs = Presentation()
    
    # 设置 16:9 页面比例
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # 核心色彩定义
    COLOR_BG = RGBColor(22, 26, 30)          # 极夜深蓝背景
    COLOR_CARD = RGBColor(38, 42, 46)        # 深灰色卡片背景
    COLOR_TEXT_MAIN = RGBColor(255, 255, 255) # 纯白主字
    COLOR_TEXT_MUTED = RGBColor(165, 175, 185) # 灰蓝色次要文字
    COLOR_ACCENT_BLUE = RGBColor(0, 210, 255) # 亮天蓝强调色
    COLOR_ACCENT_CYAN = RGBColor(0, 255, 180) # 高亮青色辅助强调
    COLOR_BORDER_GREY = RGBColor(60, 66, 74) # 卡片暗灰色边框
    
    # 字体定义
    FONT_TITLE = "Trebuchet MS"
    FONT_BODY = "Calibri"

    # =========================================================================
    # 布局辅助函数
    # =========================================================================
    
    def apply_bg(slide):
        """为幻灯片填充深邃背景"""
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background() # 无边框
        return bg

    def draw_banner_title(slide, text):
        """画页面顶部的修饰横幅和页面标题"""
        # 顶部左侧的高亮青色装饰小方块
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(0.15), Inches(0.5))
        badge.fill.solid()
        badge.fill.fore_color.rgb = COLOR_ACCENT_CYAN
        badge.line.fill.background()
        
        # 页面标题
        t_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.4), Inches(11.0), Inches(0.7))
        tf = t_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = FONT_TITLE
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEXT_MAIN

    def draw_styled_card(slide, left, top, width, height, border_color=COLOR_BORDER_GREY, fill_color=COLOR_CARD):
        """画一块带边框的高级卡片背景"""
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = fill_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
        return card

    def draw_image_placeholder(slide, left, top, width, height, title_text, instruction_text, is_uml=True):
        """画一个专门放置 UML 图、架构图或软件截图的虚线艺术化占位容器"""
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(28, 32, 38)
        frame.line.color.rgb = COLOR_ACCENT_BLUE
        frame.line.width = Pt(1.5)
        
        box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_icon = tf.paragraphs[0]
        p_icon.text = "📊 [ 架构图占位区 ]" if is_uml else "💻 [ 软件演示截图占位区 ]"
        p_icon.alignment = PP_ALIGN.CENTER
        p_icon.font.name = FONT_TITLE
        p_icon.font.size = Pt(14)
        p_icon.font.bold = True
        p_icon.font.color.rgb = COLOR_ACCENT_CYAN if is_uml else COLOR_ACCENT_BLUE
        p_icon.space_after = Pt(12)
        
        p_title = tf.add_paragraph()
        p_title.text = title_text
        p_title.alignment = PP_ALIGN.CENTER
        p_title.font.name = FONT_TITLE
        p_title.font.size = Pt(13)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_MAIN
        p_title.space_after = Pt(8)
        
        p_inst = tf.add_paragraph()
        p_inst.text = instruction_text
        p_inst.alignment = PP_ALIGN.CENTER
        p_inst.font.name = FONT_BODY
        p_inst.font.size = Pt(10.5)
        p_inst.font.color.rgb = COLOR_TEXT_MUTED
        
        return frame

    # =========================================================================
    # SLIDE 1: 封面页 (非对称科技感重构)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_bg(slide1)
    
    # 左侧装饰粗立柱 (青色到天蓝的视觉条)
    bar = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.8), Inches(0.2), Inches(4.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT_CYAN
    bar.line.fill.background()
    
    # 标题文本框
    title_box = slide1.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(11.0), Inches(3.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "高性能分布式即时通讯系统的\n体系结构设计与质量评估报告"
    p1.font.name = FONT_TITLE
    p1.font.size = Pt(38)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_ACCENT_BLUE
    p1.space_after = Pt(16)
    
    p2 = tf.add_paragraph()
    p2.text = "以自研 Simple_IMChat 架构实现为例"
    p2.font.name = FONT_BODY
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLOR_TEXT_MAIN
    
    # 底部署名与修饰线
    line = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(5.2), Inches(5.0), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_BORDER_GREY
    line.line.fill.background()
    
    info_box = slide1.shapes.add_textbox(Inches(1.3), Inches(5.4), Inches(10.0), Inches(1.0))
    itf = info_box.text_frame
    ip = itf.paragraphs[0]
    ip.text = "汇报人: toh  |  指导老师: 赵庆玲 / 李晅松  |  南京理工大学计算机学院"
    ip.font.name = FONT_BODY
    ip.font.size = Pt(15)
    ip.font.color.rgb = COLOR_TEXT_MUTED
    
    slide1.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 1 封面】\n"
        "各位老师、同学们，大家晚上好！我是汇报人 toh。今天我分享的大作业题目是"
        "《高性能分布式即时通讯系统的体系结构设计与评估》。我们将结合本学期所学的"
        "软件体系结构核心理论，深入剖析我做的项目 Simple_IMChat 的异构混合风格、"
        "4+1 视图建模，并使用效用树对该系统的动态运行外部质量与开发态静态内部质量"
        "进行具体的场景评估与设计权衡。下面请允许我开始汇报。"
    )

    # =========================================================================
    # SLIDE 2: 选题背景与高并发痛点 (双容器左右卡片对比)
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide2)
    draw_banner_title(slide2, "选题背景与即时通讯系统设计的核心挑战")
    
    # 左栏卡片：痛点
    draw_styled_card(slide2, Inches(0.8), Inches(1.6), Inches(5.6), Inches(4.8), border_color=RGBColor(200, 50, 50))
    left_box = slide2.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.2), Inches(4.4))
    ltf = left_box.text_frame
    ltf.word_wrap = True
    
    lp_title = ltf.paragraphs[0]
    lp_title.text = "🔴 即时通讯 (IM) 设计核心痛点："
    lp_title.font.size = Pt(19)
    lp_title.font.bold = True
    lp_title.font.color.rgb = RGBColor(255, 90, 90)
    lp_title.space_after = Pt(16)
    
    bullets = [
        "[ 01 ] 高并发长连接：单机需要维持并分发上万个 TCP 活跃长连接事件，极易发生网络阻塞。",
        "[ 02 ] 网络抖动丢包：移动网络下长连接瞬断，传统协议感知慢，引发丢包和时序乱序。",
        "[ 03 ] 水平伸缩性差：单体业务节点深度耦合，多机部署难以中转和同步路由状态。",
        "[ 04 ] 数据库自愈弱：依赖容器冷启动时若时序颠倒连接数据库失败，极易导致系统崩溃雪崩。"
    ]
    for b in bullets:
        bp = ltf.add_paragraph()
        bp.text = b
        bp.font.size = Pt(13)
        bp.font.color.rgb = COLOR_TEXT_MAIN
        bp.space_after = Pt(12)
        
    # 右栏卡片：策略
    draw_styled_card(slide2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8), border_color=COLOR_ACCENT_CYAN)
    right_box = slide2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "🔵 Simple_IMChat 架构应对策略："
    rp_title.font.size = Pt(19)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_CYAN
    rp_title.space_after = Pt(16)
    
    sols = [
        "[ 01 ] Epoll Reactor 网络层：底层引入网络多路复用，配合多工作线程池进行并发派发。",
        "[ 02 ] 双向 ACK 确认机制：应用层引入心跳超时重传和滑动去重窗口，保障消息 At-Least-Once 送达。",
        "[ 03 ] 无状态分布式寻址：各节点向 Redis 状态网关动态挂载路由，跨服通过 Pub/Sub 解耦。",
        "[ 04 ] Lazy Reconnect 懒重连：在业务层获取连接时懒惰检测，执行毫秒级重连自愈。"
    ]
    for s in sols:
        sp = rtf.add_paragraph()
        sp.text = s
        sp.font.size = Pt(13)
        sp.font.color.rgb = COLOR_TEXT_MAIN
        sp.space_after = Pt(12)

    slide2.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 2 选题背景】\n"
        "在架构设计之初，我们需要理清 IM 系统的核心痛点。对于长连接即时通讯系统来说，"
        "有四大物理限制：第一是高并发连接 the 性能瓶颈；第二是网络抖动下的可靠性变差；"
        "第三是单点故障与不易横向扩容；第四是由于依赖项未就绪导致的冷启动失败。"
        "我们的 Simple_IMChat 针对这四个痛点，分别设计了 Epoll Reactor 多路复用、"
        "应用层双向 ACK 投递确认、基于 Redis 的无状态状态网关、以及懒重连（Lazy Reconnect）"
        "自愈连接池。这些设计决策将在报告后面详细分析。"
    )

    # =========================================================================
    # SLIDE 3: 混合体系结构风格设计 (图+文左右卡片)
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_bg(slide3)
    draw_banner_title(slide3, "异构混合体系结构风格的选择 (Heterogeneous Style)")
    
    # 左侧：UML/风格架构图占位
    draw_image_placeholder(
        slide3, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 3.1 IMChat 异构混合风格总架构拓扑图",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 1 张：总架构拓扑图”\n截图粘贴在此处。包含客户端、Nginx、双节点和存储层。"
    )
    
    # 右侧：4个混合风格解析卡片 (纵向排列)
    y_starts = [1.6, 2.8, 4.0, 5.2]
    styles = [
        ("C/S 客户机/服务器风格", "划分系统级物理计算。胖客户端 (Qt) 承载 UI 与 ACK 发送队列，ChatServer 专注于高并发连接与路由分发。"),
        ("分层体系结构风格 (Layered)", "规范内部模块依赖。客户端分三层，服务端切为网络驱动、业务分发、实体层与 DAL 连接池四层。"),
        ("事件/消息驱动风格 (Reactor)", "高并发与水平扩展。单机内部基于 Epoll Reactor 驱动，跨服时通过 Redis Pub/Sub 异步广播解耦。"),
        ("共享中心仓库风格 (Repository)", "数据同步核心。Redis 哈希结构 user:route 存放动态路由状态，异步写盘线程批量落盘 MySQL。")
    ]
    
    for i, (name, val) in enumerate(styles):
        cur_y = Inches(y_starts[i])
        draw_styled_card(slide3, Inches(6.9), cur_y, Inches(5.6), Inches(1.15))
        
        box = slide3.shapes.add_textbox(Inches(7.1), cur_y + Inches(0.1), Inches(5.2), Inches(0.95))
        tf = box.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"[ {i+1:02d} ] {name}"
        p.font.name = FONT_TITLE
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT_BLUE
        p.space_after = Pt(4)
        
        pd = tf.add_paragraph()
        pd.text = val
        pd.font.name = FONT_BODY
        pd.font.size = Pt(11)
        pd.font.color.rgb = COLOR_TEXT_MAIN
        
    slide3.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 3 混合架构风格】\n"
        "我们系统的核心设计在于它不是单一风格，而是一个异构混合体系结构风格。\n"
        "在幻灯片左侧的架构拓扑图里，老师可以看到，它在物理大格局上是经典 C/S 风格，客户端承载了大量的 UI "
        "以及双向可靠 ACK 机制。在纵向依赖上，我们严格采用分层风格将各层间调用解耦。单机的高吞吐基于 Linux "
        "Epoll Reactor 事件就绪触发回调；而在多机集群下，我们通过 Redis 发布订阅的消息总线抹平了物理实例差异。"
        "最后，以 Redis 作为在线状态共享网关和 MySQL Timeline 历史大仓库，体现了中心仓库风格。"
    )

    # =========================================================================
    # SLIDE 4: 4+1 视图之 场景视图
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_bg(slide4)
    draw_banner_title(slide4, "“4+1” 视图模型之 场景视图 (Scenarios / Use-Case View)")
    
    # 左侧：场景用例图占位
    draw_image_placeholder(
        slide4, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 4.1 Simple_IMChat 社交功能用例模型",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 2 张：场景用例图”\n截图粘贴在此处。包含登录、单聊、群聊、离线拉取等用例。"
    )
    
    # 右侧：5个核心用例场景解析 (卡片布局)
    draw_styled_card(slide4, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide4.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "系统核心用例与场景划分："
    rp_title.font.size = Pt(18)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(14)
    
    usecases = [
        "[ 01 ] 注册与安全登录：提供基于 Session Token 的状态校验与加盐 MD5 密码认证。",
        "[ 02 ] 添加好友与实时状态感知：支持双向好友申请审批流，用户上线/下线通过状态网关实时推送。",
        "[ 03 ] 单对单即时聊天：最核心的长连接交互，使用应用层序列号与待确认队列保障投递。",
        "[ 04 ] 群组创建与群发：支持多播逻辑。由业务层解析群成员列表，执行跨节点的广播中转。",
        "[ 05 ] Timeline 离线同步：增量式历史同步，用户上线时仅拉取缺失的消息 ID 区间。"
    ]
    for uc in usecases:
        up = rtf.add_paragraph()
        up.text = uc
        up.font.size = Pt(12)
        up.font.color.rgb = COLOR_TEXT_MAIN
        up.space_after = Pt(10)

    slide4.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 4 场景视图】\n"
        "“4+1”视图模型以场景视图为核心驱动。如图 4.1 所示，我们的即时通讯用例包含了"
        "登录鉴权、好友申请与状态感知、单聊、群聊、以及增量 Timeline 离线拉取。\n"
        "这五个用例覆盖了即时通讯中最具代表性的功能与非功能边界，所有的体系结构决策（如 Reactor、"
        "双向 ACK 和分布式路由）都是为了支撑这五个用例在高并发环境下的稳定运行。"
    )

    # =========================================================================
    # SLIDE 5: 4+1 视图之 逻辑视图
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_bg(slide5)
    draw_banner_title(slide5, "“4+1” 视图模型之 逻辑视图 (面向对象类模型)")
    
    # 左侧：类图占位
    draw_image_placeholder(
        slide5, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 5.1 客户端 (IMClient) 与服务端 (ChatService) 类设计图",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 3 张：客户端类图” 与 “第 4 张：服务端类图”\n截图拼接并粘贴在此处。展示主回调和映射分发。"
    )
    
    # 右侧：类职责卡片
    draw_styled_card(slide5, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide5.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "核心类的设计职责与映射机制："
    rp_title.font.size = Pt(18)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(14)
    
    classes = [
        "🔵 客户端 IMClient (imclient.cpp)：\n作为网络与控制核心，封装 QTcpSocket。维护 pendingSendAckMap_（待确认发送队列）以驱动重传，处理粘包分包，并提供统一的业务槽函数接口。",
        "🔵 服务端 ChatService (chatservice.cpp)：\n采用单例模式。内部维护成员变量 _msgHandlerMap 消息处理器映射表。当收到网络包时，根据其 MsgType（如单聊、创建群）直接索引回调 Handler，消除了硬编码耦合。"
    ]
    for c in classes:
        cp = rtf.add_paragraph()
        cp.text = c
        cp.font.size = Pt(12)
        cp.font.color.rgb = COLOR_TEXT_MAIN
        cp.space_after = Pt(14)

    slide5.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 5 逻辑视图】\n"
        "逻辑视图关注系统的静态结构。在面向对象设计中，我们重点设计了客户端的 IMClient 类和"
        "服务端的 ChatService 类。如右侧职责所示，IMClient 封装了套接字并管理重传状态机；"
        "而 ChatService 作为服务端的单例，核心是维护了一个消息回调映射表。这种设计彻底摆脱了"
        "传统 switch-case 分发消息带来的高耦合，使系统具备了极佳的可维护性。"
    )

    # =========================================================================
    # SLIDE 6: 4+1 视图之 过程视图 (一) - 可靠投递
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_bg(slide6)
    draw_banner_title(slide6, "过程视图 (Process View) - 应用层双向 ACK 可靠消息投递")
    
    # 左侧：双向 ACK 时序图占位
    draw_image_placeholder(
        slide6, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 6.1 应用层双向 ACK 与滑动去重窗口时序流程图",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 5 张：应用层双向 ACK 时序图”\n截图粘贴在此处。展现发送待确认与服务端重传校验。"
    )
    
    # 右侧：过程步骤卡片
    draw_styled_card(slide6, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide6.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "可靠投递的过程控制时序："
    rp_title.font.size = Pt(18)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(16)
    
    steps = [
        "[ 01 ] 客户端发送缓存：发送端生成自增 seq，写入 pendingSendAckMap_，同时在后台挂载 3 秒重传定时器。",
        "[ 02 ] 服务端滑动窗口校验：服务端收到数据包后先执行去重检测，若是由于重传引起的重复消息则直接丢弃，但仍强制回送确认 ACK。",
        "[ 03 ] 接收端消费并响应：接收端安全将消息写入本地，在 UI 上成功渲染后，向服务端回复 ACK 确认数据包。",
        "[ 04 ] 发送端注销定时器：发送方收到 ACK 确认，根据 seq 移除本地缓存并销毁重传定时器。若超时（3秒）则触发自动补发。"
    ]
    for s in steps:
        sp = rtf.add_paragraph()
        sp.text = s
        sp.font.size = Pt(12)
        sp.font.color.rgb = COLOR_TEXT_MAIN
        sp.space_after = Pt(12)

    slide6.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 6 过程视图一】\n"
        "过程视图关注动态的并发和时序。这里展示了应用层双向 ACK 投递确认时序。移动网络环境下，"
        "底层 TCP 协议因其超长的重传时序无法满足即时通讯的低延迟高可用要求。因此，我们在应用层"
        "引入了双确认机制。客户端发包时将其加入 pendingSendAckMap_ 并挂载定时器；服务端接收时，"
        "通过滑动去重窗口防止由于重传导致的消息重复入库，从而保障了系统的 At-Least-Once 可靠性。"
    )

    # =========================================================================
    # SLIDE 7: 4+1 视图之 过程视图 (二) - 分布式路由
    # =========================================================================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_bg(slide7)
    draw_banner_title(slide7, "过程视图 (Process View) - 分布式多节点 Pub/Sub 路由转发")
    
    # 左侧：分布式路由时序图占位
    draw_image_placeholder(
        slide7, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 7.1 分布式跨服寻址与消息通道发布订阅流转图",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 6 张：跨节点分布式路由时序图”\n截图粘贴在此处。展现 Server1 读 Redis 状态网关并发布至 channel。"
    )
    
    # 右侧：过程步骤卡片
    draw_styled_card(slide7, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide7.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "分布式长连接路由的中转时序："
    rp_title.font.size = Pt(18)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(16)
    
    steps = [
        "[ 01 ] 路由本地检索：UserA 在 Server1 发送单聊消息，Server1 检索 _userConnMap 发现 UserB 不在本节点登录。",
        "[ 02 ] 集中式寻址：Server1 查询 Redis 共享状态哈希表 user:route。通过 HGET 获取目标在线位置为 Server2 的 IP 与端口。",
        "[ 03 ] 异步通道发布：Server1 向 Redis 通道 chat_server_channel_2 执行 PUBLISH 广播，将消息数据投递出去。",
        "[ 04 ] 订阅派发：Server2 后台监听该通道，订阅收到后，在其本地连接表中定位 UserB 的 Socket 并推送到接收端。"
    ]
    for s in steps:
        sp = rtf.add_paragraph()
        sp.text = s
        sp.font.size = Pt(12)
        sp.font.color.rgb = COLOR_TEXT_MAIN
        sp.space_after = Pt(12)

    slide7.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 7 过程视图二】\n"
        "过程视图的第二个重点是分布式长连接路由。在多物理节点下，用户连接是分散的。"
        "当发送端向另一个节点的接收端发送消息时，Server1 会检索 Redis 状态路由网关定位"
        "目标用户所在的服务器。随后通过 Redis 异步发布订阅（Pub/Sub）机制完成跨服消息中转。"
        "这使得各业务节点完全无状态，实现了平滑的水平扩展。"
    )

    # =========================================================================
    # SLIDE 8: 4+1 视图之 开发视图 (模块与依赖)
    # =========================================================================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_bg(slide8)
    draw_banner_title(slide8, "“4+1” 视图模型之 开发视图 (Development View)")
    
    # 左侧：开发依赖图占位
    draw_image_placeholder(
        slide8, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 8.1 客户端与服务端模块静态编译依赖图",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 10 张：开发视图模块编译依赖图”\n截图粘贴在此处。展示 UI层、网络层、业务层与三方库依赖。"
    )
    
    # 右侧：开发视图说明卡片
    draw_styled_card(slide8, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide8.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "静态代码模块的开发组织约束："
    rp_title.font.size = Pt(18)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(14)
    
    devs = [
        "[ 01 ] 编译环境隔离：服务端面向 Linux (Ubuntu/WSL2)，基于 CMake 和 POSIX 接口编译；客户端面向 Windows，基于 Qt 6.10/MSVC2022 构建。",
        "[ 02 ] 共享通信契约 (proto)：两端共享 Protobuf 协议文件，通过二进制编解码进行解耦，保证两端独立并行开发。",
        "[ 03 ] 分层静态依赖约束：严格遵守“上层只能依赖本层或下层”的软件体系结构规则，下层模块（如DAL连接池）对上层业务完全隐藏细节。",
        "[ 04 ] 依赖管理机制：采用 vcpkg 管理本地 Protobuf 静态库链接，防止多机编译时因 DLL 缺失发生运行时崩溃。"
    ]
    for d in devs:
        dp = rtf.add_paragraph()
        dp.text = d
        dp.font.size = Pt(11.5)
        dp.font.color.rgb = COLOR_TEXT_MAIN
        dp.space_after = Pt(8)

    slide8.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 8 开发视图】\n"
        "开发视图描述了代码模块在开发环境下的静态组织。如左侧图 8.1 所示，我们通过接口隔离"
        "将系统划分为客户端静态依赖树、服务端静态依赖树以及它们共享的 proto 协议定义。\n"
        "我们严格遵守了‘上层只能依赖本层或下层’的体系结构规范，客户端 UI 依赖网络层，"
        "网络层依赖底座；服务端启动层依赖业务层，业务层依赖模型与连接池。这种层次化的静态"
        "组织极大地提升了系统的可维护性与测试性。"
    )

    # =========================================================================
    # SLIDE 9: 4+1 视图之 物理部署视图
    # =========================================================================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_bg(slide9)
    draw_banner_title(slide9, "“4+1” 视图模型之 物理部署视图 (Physical View)")
    
    # 左侧：物理部署拓扑图占位
    draw_image_placeholder(
        slide9, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 9.1 基于 Docker Compose 的物理部署网络拓扑图",
        "【操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 8 张：Docker Compose 物理部署拓扑图”\n截图粘贴在此处。包含外网、Nginx代理、两台ChatServer及存储容器。"
    )
    
    # 右侧：部署说明卡片
    draw_styled_card(slide9, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide9.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "系统物理拓扑与容器编排："
    rp_title.font.size = Pt(18)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(16)
    
    phys = [
        "[ 01 ] Nginx 负载均衡网关 (Port: 8000)：作为七层或四层反向代理最前端，暴露公网提供 TCP 轮询分流。",
        "[ 02 ] 双业务节点容器 (Port: 8888 / 8889)：两台无状态 ChatServer 容器部署于桥接私有网段内，分担计算负荷。",
        "[ 03 ] 数据共享与缓存容器 (Port: 6379)：部署 Redis 容器，作为多节点之间的在线路由与状态网关。",
        "[ 04 ] 持久化数据库容器 (Port: 3307)：部署 MySQL 容器，映射至内部 3306，确保 Timeline 历史强一致性存盘。"
    ]
    for p in phys:
        pp = rtf.add_paragraph()
        pp.text = p
        pp.font.size = Pt(12)
        pp.font.color.rgb = COLOR_TEXT_MAIN
        pp.space_after = Pt(12)

    slide9.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 9 物理视图】\n"
        "物理视图展示了系统构件到硬件节点的映射。我们设计了一套 Docker 容器化编排网络。\n"
        "外网用户通过端口 8000 接入 Nginx 负载均衡网关，Nginx 通过轮询算法分流至后方的"
        "两台独立 ChatServer 容器（分别监听 8888 和 8889）。业务节点在私有桥接网络中"
        "与后方的 Redis 共享状态网关和 MySQL 8.0 物理库通信。这实现了一键容器化快速编排与部署。"
    )

    # =========================================================================
    # SLIDE 10: 质量属性评估 (效用树与情景)
    # =========================================================================
    slide8_2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide8_2)
    draw_banner_title(slide8_2, "软件体系结构质量评估 - 效用树与三要素情景")
    
    # 左侧：效用树模型图占位
    draw_image_placeholder(
        slide8_2, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 10.1 IMChat 质量属性效用树 (Utility Tree) 分类结构图",
        "─操作提示】：请将您从《UML及架构图大全.md》导出的\n“第 9 张：体系结构质量属性效用树”\n截图粘贴在此处。包含内部开发态质量与外部运行态质量的分级。"
    )
    
    # 右侧：说明卡片
    draw_styled_card(slide8_2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide8_2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "【激励 - 环境 - 响应】三要素情景："
    rp_title.font.size = Pt(17)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_BLUE
    rp_title.space_after = Pt(14)
    
    steps = [
        "🔵 可用性评估情景：网络切网引发物理套接字异常瞬断。\n• 激励：发生局部长连接消息收发丢包。\n• 环境：用户高频聊天的活跃运行状态下。\n• 响应：客户端心跳定时器 3 秒重传，去重窗口幂等渲染，保障At-Least-Once。",
        "🔵 性能评估情景：高峰期服务器遭遇每秒数万条消息写入压力。\n• 激励：数据需以高吞吐落盘持久化数据库 timeline_history 表中。\n• 环境：网络 I/O 与 CPU 高负载状态下。\n• 响应：主线程 Reactor 异步解耦，先写 Redis ZSet，异步存盘线程批量落盘。"
    ]
    for s in steps:
        sp = rtf.add_paragraph()
        sp.text = s
        sp.font.size = Pt(11.5)
        sp.font.color.rgb = COLOR_TEXT_MAIN
        sp.space_after = Pt(10)

    slide8_2.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 10 质量属性评估】\n"
        "按照课程课件的体系结构评估大纲要求，质量保障需区分为内部静态质量和外部动态质量。\n"
        "首先，在开发态内部质量上，我们展示了 Protobuf 带来的高可重用性、Qt 实现的平台无关可移植性、"
        "以及接口解耦实现的可测试性。\n"
        "其次，在运行态外部质量上，我们严格遵循‘激励-环境-响应’的三要素框架进行情景评估。"
        "对于可用性，我们的激励是连接瞬断，系统在活跃环境下通过心跳定时器重传和去重窗口来响应；"
        "对于性能，我们的激励是高频写盘，主线程 Reactor 避免同步写盘，而通过写 Redis 并使用异步落盘线程池"
        "批量刷写 MySQL 来响应。这为系统在大并发下的表现提供了定性评估。"
    )

    # =========================================================================
    # SLIDE 11: 架构设计三大核心亮点与设计折中 (三栏卡片)
    # =========================================================================
    slide9_2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide9_2)
    draw_banner_title(slide9_2, "Simple_IMChat 架构优化亮点与设计决策折中")
    
    # 绘制三个亮点卡片 (横向网格)
    card_width = Inches(3.7)
    card_height = Inches(4.8)
    card_y = Inches(1.6)
    
    highlights = [
        ("亮点 1. 双向 ACK 与滑动去重", 
         "🔴 敏感点 (可靠性提升)：\n"
         "客户端待确认队列 (pendingSendAckMap_) 结合超时重传，解决网络瞬断丢包；滑动窗口去重保障业务幂等渲染。\n\n"
         "🔵 设计折中 (Trade-off)：\n"
         "客户端增加了本地内存缓存开销；心跳定时检测增加了约 15% 协议小包网络带宽流量。"),
         
        ("亮点 2. Lazy Reconnect 数据库池自愈", 
         "🔴 敏感点 (抗抖动可用性)：\n"
         "自研连接池，在业务线程调用 getConnection() 获取连接时，前置懒惰执行 mysql_ping() 检测，连接断开时秒级后台重连。\n\n"
         "🔵 架构收益 (高容错)：\n"
         "服务具备极强时序防抖能力。Docker compose 拉起容器时即便数据库 mysql-chat 初始化未就绪，服务端亦可冷启动防崩溃自愈。"),
         
        ("亮点 3. Redis 缓存与异步批量存盘", 
         "🔴 敏感点 (高并发与吞吐量)：\n"
         "主线程网络 Epoll Reactor 仅负责网络事件就绪，将计算投递给工作线程。写 Redis 缓存快速返回回复。\n\n"
         "🔵 设计折中 (数据一致性)：\n"
         "由异步写盘线程池定时刷入 MySQL Timeline 历史消息表 (timeline_history)。牺牲微秒级数据强一致性，换取万级超低延迟吞吐。")
    ]
    
    for i, (name, desc) in enumerate(highlights):
        card_x = Inches(0.8 + i * 4.0)
        draw_styled_card(slide9_2, card_x, card_y, card_width, card_height)
        
        box = slide9_2.shapes.add_textbox(card_x + Inches(0.1), card_y + Inches(0.15), card_width - Inches(0.2), card_height - Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = True
        
        p_name = tf.paragraphs[0]
        p_name.text = f"[ 亮点 {i+1} ]"
        p_name.font.name = FONT_TITLE
        p_name.font.size = Pt(13)
        p_name.font.color.rgb = COLOR_ACCENT_CYAN
        
        p_sub = tf.add_paragraph()
        p_sub.text = name.split('\n')[0]
        p_sub.font.name = FONT_TITLE
        p_sub.font.size = Pt(16)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_TEXT_MAIN
        p_sub.space_after = Pt(14)
        
        p_desc = tf.add_paragraph()
        p_desc.text = desc
        p_desc.font.name = FONT_BODY
        p_desc.font.size = Pt(12)
        p_desc.font.color.rgb = COLOR_TEXT_MAIN
        
    slide9_2.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 11 架构亮点与折中】\n"
        "这里我向老师总结我们在 Simple_IMChat 架构上付出的优化努力和设计折中。\n"
        "第一，应用层双向 ACK 保障了网络瞬断下的消息可靠投递。它的设计折中是稍微增加了心跳包"
        "网络带宽开销，但换取了极高的数据可用性。\n"
        "第二，是 Lazy Reconnect 数据库自愈机制，它成功解决了 Docker compose 部署时的"
        "服务冷启动时序问题，服务秒级防雪崩自愈。\n"
        "第三，是 Redis ZSet 缓存加异步存盘线程池。主线程 Reactor 和慢速磁盘 IO 彻底隔离，"
        "它的设计折中是在微秒级极极端崩溃下，可能牺牲少量的 Timeline 强一致性，但为我们带来了"
        "万级高吞吐的并发性能提升。这是一项经典的性能与可用性折中决策。"
    )

    # =========================================================================
    # SLIDE 12: 软件运行演示展示页 (演示截图占位 & 答辩Q&A)
    # =========================================================================
    slide10_2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide10_2)
    draw_banner_title(slide10_2, "Simple_IMChat 软件运行演示与答辩准备")
    
    # 左侧：软件截图展示占位
    draw_image_placeholder(
        slide10_2, 
        Inches(0.8), Inches(1.6), Inches(5.8), Inches(4.8),
        "图 12.1 Simple_IMChat 客户端登录及文字交互演示图",
        "【操作提示】：请在此处贴上客户端的物理运行截图：\n"
        "1. 客户端登录界面（输入用户 ID 和密码登录）。\n"
        "2. 好友列表主窗体（查看实时上线的绿色在线指示）。\n"
        "3. 核心聊天对话框（展现两台客户端正在通过网关进行即时单聊发包）。",
        is_uml=False
    )
    
    # 右栏卡片：Q&A
    draw_styled_card(slide10_2, Inches(6.9), Inches(1.6), Inches(5.6), Inches(4.8))
    right_box = slide10_2.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.2), Inches(4.4))
    rtf = right_box.text_frame
    rtf.word_wrap = True
    
    rp_title = rtf.paragraphs[0]
    rp_title.text = "❓ 答辩可能质疑点与应答技巧"
    rp_title.font.size = Pt(19)
    rp_title.font.bold = True
    rp_title.font.color.rgb = COLOR_ACCENT_CYAN
    rp_title.space_after = Pt(14)
    
    sols = [
        "[ 01 ] 问：若中心化 Redis 发生故障路由崩溃怎么响应？\n答：物理视图中已支持部署 Redis 哨兵集群。业务层有 Dalton 容错连接，挂掉路由后自动降级为单节点通信，将连接状态移入本地缓存。",
        "[ 02 ] 问：应用层双向 ACK 会不会引起网络拥塞与高包开销？\n答：心跳采用指数退避算法，且通过滑动去重窗口去除了不必要的重复确认。外网 Nginx L4 负载均衡对请求轮询分流，整体表现极低延迟。"
    ]
    for s in sols:
        sp = rtf.add_paragraph()
        sp.text = s
        sp.font.size = Pt(12)
        sp.font.color.rgb = COLOR_TEXT_MAIN
        sp.space_after = Pt(10)

    slide10_2.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 12 总结与展望】\n"
        "最后，针对系统整体体系结构，我提出两个未来演进方向：\n"
        "一是将单体 ChatServer 演进为微服务，剥离出鉴权、群组等逻辑；"
        "二是引入专业的 Kafka 消息中间件来应对百万级并发路由。\n"
        "在答辩准备中，我们对 Redis 状态网关单点故障、双向 ACK 的小包网络负荷等设计细节进行了容错与降级演进规划，"
        "这确保了整个软件体系结构的演进灵活性。\n"
        "以上就是我关于 Simple_IMChat 软件体系结构设计大作业的完整汇报。感谢各位老师的聆听，请老师批评指正！"
    )

    # =========================================================================
    # SLIDE 13: 结束致谢页 (Thank You)
    # =========================================================================
    slide11_2 = prs.slides.add_slide(blank_layout)
    apply_bg(slide11_2)
    
    # 中央致谢卡片
    draw_styled_card(slide11_2, Inches(2.5), Inches(2.0), Inches(8.33), Inches(3.5), border_color=COLOR_ACCENT_CYAN)
    
    thanks_box = slide11_2.shapes.add_textbox(Inches(2.7), Inches(2.3), Inches(7.93), Inches(2.9))
    ttf = thanks_box.text_frame
    ttf.word_wrap = True
    
    tp1 = ttf.paragraphs[0]
    tp1.alignment = PP_ALIGN.CENTER
    tp1.text = "THANK YOU"
    tp1.font.name = FONT_TITLE
    tp1.font.size = Pt(54)
    tp1.font.bold = True
    tp1.font.color.rgb = COLOR_ACCENT_BLUE
    tp1.space_after = Pt(14)
    
    tp2 = ttf.add_paragraph()
    tp2.alignment = PP_ALIGN.CENTER
    tp2.text = "感谢各位老师的聆听与指导！"
    tp2.font.name = FONT_BODY
    tp2.font.size = Pt(22)
    tp2.font.color.rgb = COLOR_TEXT_MAIN
    
    slide11_2.notes_slide.notes_text_frame.text = (
        "【现场演讲稿 - Slide 13 致谢】\n"
        "再次感谢各位老师！以上就是我的课设答辩汇报。请老师提问。"
    )

    # 物理保存文件为 v4
    output_path = "Simple_IMChat_汇报幻灯片_v4.pptx"
    prs.save(output_path)
    print(f"PPTX 物理文件已成功优化并输出至: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_deck()
